"""PPTX extraction via python-pptx.

Each slide's text shapes are joined; speaker notes are appended after a
"Notes:" prefix when present. Slide order is preserved.
"""

from __future__ import annotations

from . import ExtractionError, ExtractionResult


def extract(pptx_path: str) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ExtractionError(f"python-pptx not installed: {exc}") from exc

    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        raise ExtractionError(f"Could not open PPTX: {exc}") from exc

    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_chunks: list[str] = [f"# Slide {i}"]

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_chunks.append(text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_chunks.append(f"Notes: {notes}")

        if len(slide_chunks) > 1:
            slides.append("\n".join(slide_chunks))

    text = "\n\n".join(slides)
    if not text.strip():
        raise ExtractionError("PPTX appears to be empty")

    return ExtractionResult(text=text, page_count=len(prs.slides))
